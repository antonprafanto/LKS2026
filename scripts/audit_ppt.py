"""Audit briefing PowerPoint for readability and content issues."""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

PPTX = Path(__file__).resolve().parent.parent / "docs" / "Briefing-Juri-LKS2026-Robot-Otonom.pptx"
FOOTER_Y = 7.05
CONTENT_MAX_Y = 6.85

BUGS: list[str] = []
WARNS: list[str] = []


def bug(msg: str) -> None:
    BUGS.append(msg)


def warn(msg: str) -> None:
    WARNS.append(msg)


def shape_bottom(shape) -> float:
    return shape.top.inches + shape.height.inches


def main() -> int:
    prs = Presentation(PPTX)
    all_text = []

    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    slide_text.append(t)
                    all_text.append(t)

            bottom = shape_bottom(shape)
            if bottom > CONTENT_MAX_Y and shape.top.inches < FOOTER_Y - 0.05:
                if not (shape.top.inches >= FOOTER_Y - 0.01):
                    warn(f"Slide {i}: konten y={bottom:.2f} mendekati footer")

        joined = " ".join(slide_text)
        if len(joined) > 500:
            warn(f"Slide {i}: teks panjang ({len(joined)} karakter) — pertimbangkan pecah slide")
        if len(slide_text) == 1 and "Demo Modul E" in joined and len(joined) < 120:
            pass  # section ok

    corpus = "\n".join(all_text)
    if ".csv" in corpus.lower():
        bug("Masih ada referensi file CSV (sudah diganti Excel)")
    if "6,67" in corpus or "6.67" in corpus:
        bug("Masih menyebut bobot 6,67 (seharusnya 60/9 di Excel)")
    if "lembar-undian-kubus.csv" in corpus:
        bug("Referensi lembar-undian-kubus.csv usang")
    if "baris 22" not in corpus.lower():
        warn("PPT belum menjelaskan isian marker baris 22")
    if "CONTOH" not in corpus:
        warn("PPT belum menyebut file LKS2026-Penilaian-Juri-CONTOH.xlsx")

    if len(prs.slides) < 18:
        warn(f"Hanya {len(prs.slides)} slide — mungkin kurang lengkap")

    print(f"=== AUDIT PPT: {PPTX.name} ({len(prs.slides)} slide) ===")
    for b in BUGS:
        print("[BUG]", b)
    for w in WARNS:
        print("[WARN]", w)
    if not BUGS and not WARNS:
        print("Tidak ada masalah terdeteksi.")
    return 1 if BUGS else 0


if __name__ == "__main__":
    raise SystemExit(main())
