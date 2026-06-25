"""Generate readable branded briefing PowerPoint for LKS 2026 judges."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "assets"
OUT = ROOT / "docs" / "Briefing-Juri-LKS2026-Robot-Otonom.pptx"

# Layout constants (inches) — 10 x 7.5 slide
HEADER_H = 0.92
ACCENT_H = 0.05
FOOTER_Y = 7.02
CONTENT_TOP = 1.05
CONTENT_W = 9.0
CONTENT_LEFT = 0.5
CONTENT_H = 5.75  # aman di atas footer (y=7.02)

NAVY = RGBColor(0x1F, 0x4E, 0x79)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF7, 0xF9, 0xFC)
WARN_BG = RGBColor(0xFF, 0xF3, 0xE0)

LOGO_KEMEN = ASSETS / "logo-kemendikdasmen.png"
LOGO_LKS = ASSETS / "logo-lks-dikmen.png"
LOGO_GEO = ASSETS / "logo-lks-geometric.png"

_slide_no = 0


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_para(para, *, size=16, bold=False, color=DARK, align=None, space_after=6):
    para.font.size = Pt(size)
    para.font.bold = bold
    para.font.color.rgb = color
    if align is not None:
        para.alignment = align
    para.space_after = Pt(space_after)


def _add_branding(slide, prs, title: str = "") -> None:
    global _slide_no
    _slide_no += 1
    sw = prs.slide_width

    bar = slide.shapes.add_shape(1, 0, 0, sw, Inches(HEADER_H))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    accent = slide.shapes.add_shape(1, 0, Inches(HEADER_H), sw, Inches(ACCENT_H))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()

    if LOGO_KEMEN.exists():
        slide.shapes.add_picture(str(LOGO_KEMEN), Inches(0.12), Inches(0.1), height=Inches(0.58))
    if LOGO_LKS.exists():
        slide.shapes.add_picture(str(LOGO_LKS), Inches(8.72), Inches(0.06), height=Inches(0.68))

    if title:
        tb = slide.shapes.add_textbox(Inches(2.35), Inches(0.14), Inches(6.0), Inches(0.62))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        _set_para(p, size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    foot = slide.shapes.add_shape(1, 0, Inches(FOOTER_Y), sw, Inches(0.4))
    foot.fill.solid()
    foot.fill.fore_color.rgb = GREEN
    foot.line.fill.background()

    ft = slide.shapes.add_textbox(Inches(0.25), Inches(FOOTER_Y + 0.04), Inches(9.5), Inches(0.32))
    fp = ft.text_frame.paragraphs[0]
    fp.text = f"Slide {_slide_no}  |  Briefing Juri LKS 2026 — Robot Bergerak Otonom  |  Kemendikdasmen"
    _set_para(fp, size=9, color=WHITE, align=PP_ALIGN.CENTER, space_after=0)


def _body_box(slide, top=CONTENT_TOP, height=None):
    h = height or CONTENT_H
    box = slide.shapes.add_textbox(
        Inches(CONTENT_LEFT), Inches(top), Inches(CONTENT_W), Inches(h)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return tf


def _add_bullets(prs, title: str, bullets: list[str], *, sub: str = "", size=17):
    slide = _blank(prs)
    _add_branding(slide, prs, title)
    tf = _body_box(slide)
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        _set_para(p, size=size, space_after=10)
    if sub:
        p = tf.add_paragraph()
        p.text = sub
        _set_para(p, size=12, color=ORANGE, space_after=0)


def _add_section(prs, title: str, subtitle: str = ""):
    slide = _blank(prs)
    _add_branding(slide, prs)
    panel = slide.shapes.add_shape(
        1, Inches(0.7), Inches(2.5), Inches(8.6), Inches(2.0 if subtitle else 1.5)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = NAVY
    panel.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.75), Inches(8.0), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    _set_para(p, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        _set_para(p2, size=16, color=ORANGE, align=PP_ALIGN.CENTER, space_after=0)


def _add_table(prs, title: str, headers: list[str], rows: list[list[str]]):
    slide = _blank(prs)
    _add_branding(slide, prs, title)
    nrows, ncols = len(rows) + 1, len(headers)
    tbl_h = min(Inches(0.42 * nrows + 0.25), Inches(CONTENT_H - 0.1))
    table = slide.shapes.add_table(
        nrows, ncols, Inches(0.4), Inches(CONTENT_TOP), Inches(9.2), tbl_h
    ).table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            _set_para(p, size=12, bold=True, color=WHITE)
    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            for p in table.cell(r, c).text_frame.paragraphs:
                p.text = val
                _set_para(p, size=11)


def _add_flow(prs):
    slide = _blank(prs)
    _add_branding(slide, prs, "Alur Run Otonom (Modul E)")
    steps = [
        "1. Briefing layout arena",
        "2. Programming (laptop boleh disentuh)",
        "3. Peserta nyatakan SIAP",
        "4. Laptop dilarang | Juri acak kubus",
        "5. START → tekan tombol di robot",
        "6. Run otonom — catat di Excel",
        "7. Nilai 1/0 per kubus di Modul E",
    ]
    tf = _body_box(slide)
    for i, s in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = s
        _set_para(p, size=18, bold=("SIAP" in s or "START" in s), space_after=8)

    nb = slide.shapes.add_shape(1, Inches(0.45), Inches(5.85), Inches(9.1), Inches(0.85))
    nb.fill.solid()
    nb.fill.fore_color.rgb = WARN_BG
    nb.line.color.rgb = ORANGE
    nt = slide.shapes.add_textbox(Inches(0.6), Inches(5.95), Inches(8.8), Inches(0.7))
    np = nt.text_frame.paragraphs[0]
    np.text = (
        "PENTING: Sentuh laptop setelah SIAP atau sentuh robot setelah START "
        "= pelanggaran (run dapat dibatalkan)."
    )
    _set_para(np, size=13, bold=True, color=ORANGE, space_after=0)


def _add_demo_step(prs, step_no: int, title: str, bullets: list[str], note: str = ""):
    slide = _blank(prs)
    _add_branding(slide, prs, f"Demo Modul E — Langkah {step_no}/8")

    badge = slide.shapes.add_shape(1, Inches(0.45), Inches(CONTENT_TOP), Inches(0.55), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ORANGE
    badge.line.fill.background()
    bt = slide.shapes.add_textbox(Inches(0.45), Inches(CONTENT_TOP + 0.06), Inches(0.55), Inches(0.45))
    bp = bt.text_frame.paragraphs[0]
    bp.text = str(step_no)
    _set_para(bp, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, space_after=0)

    tt = slide.shapes.add_textbox(Inches(1.15), Inches(CONTENT_TOP - 0.02), Inches(8.3), Inches(0.55))
    tp = tt.text_frame.paragraphs[0]
    tp.text = title
    _set_para(tp, size=22, bold=True, color=NAVY, space_after=0)

    body_top = CONTENT_TOP + 0.65
    body_h = 4.85 if note else CONTENT_H - 0.65
    tf = _body_box(slide, top=body_top, height=body_h)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        _set_para(p, size=17, space_after=12)

    if note:
        nb = slide.shapes.add_shape(1, Inches(0.45), Inches(6.05), Inches(9.1), Inches(0.78))
        nb.fill.solid()
        nb.fill.fore_color.rgb = WARN_BG
        nb.line.color.rgb = ORANGE
        nt = slide.shapes.add_textbox(Inches(0.6), Inches(6.15), Inches(8.8), Inches(0.62))
        np = nt.text_frame.paragraphs[0]
        np.text = f"PENTING — Juri: {note}"
        _set_para(np, size=12, bold=True, color=ORANGE, space_after=0)


def _add_title(prs: Presentation):
    slide = _blank(prs)
    global _slide_no
    _slide_no += 1

    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(HEADER_H + ACCENT_H))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    panel = slide.shapes.add_shape(1, 0, Inches(1.1), prs.slide_width, Inches(5.6))
    panel.fill.solid()
    panel.fill.fore_color.rgb = LIGHT
    panel.line.fill.background()

    if LOGO_KEMEN.exists():
        slide.shapes.add_picture(str(LOGO_KEMEN), Inches(0.4), Inches(0.12), height=Inches(0.65))
    if LOGO_LKS.exists():
        slide.shapes.add_picture(str(LOGO_LKS), Inches(3.9), Inches(1.45), height=Inches(1.55))
    if LOGO_GEO.exists():
        slide.shapes.add_picture(str(LOGO_GEO), Inches(7.0), Inches(1.5), height=Inches(1.45))

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.35), Inches(9.0), Inches(0.9))
    p = tb.text_frame.paragraphs[0]
    p.text = "Briefing Juri"
    _set_para(p, size=38, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.25), Inches(9.0), Inches(1.6))
    tf = tb2.text_frame
    tf.word_wrap = True
    lines = [
        ("Robot Bergerak Otonom", 20, DARK, True),
        ("LKS Nasional XXXIV 2026", 17, ORANGE, False),
        ("Kemendikdasmen — Pusat Prestasi Nasional", 14, MUTED, False),
    ]
    for i, (txt, sz, col, b) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = txt
        _set_para(para, size=sz, bold=b, color=col, align=PP_ALIGN.CENTER, space_after=6)

    foot = slide.shapes.add_shape(1, 0, Inches(FOOTER_Y), prs.slide_width, Inches(0.4))
    foot.fill.solid()
    foot.fill.fore_color.rgb = GREEN
    ft = slide.shapes.add_textbox(Inches(0.25), Inches(FOOTER_Y + 0.04), Inches(9.5), Inches(0.32))
    fp = ft.text_frame.paragraphs[0]
    fp.text = "Slide 1  |  Briefing Juri LKS 2026  |  Autonomous Mobile Robotics"
    _set_para(fp, size=9, color=WHITE, align=PP_ALIGN.CENTER, space_after=0)


def main() -> None:
    if not LOGO_KEMEN.exists():
        import subprocess
        subprocess.run(["python", str(ROOT / "scripts" / "extract_branding_assets.py")], check=True)

    global _slide_no
    _slide_no = 0
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _add_title(prs)

    _add_bullets(prs, "Agenda", [
        "Peran juri & file penilaian (Excel)",
        "Modul A–E dan bobot skor",
        "Sistem penilaian: Judgement vs Measurement",
        "Alur run otonom (SIAP → START → nilai)",
        "Demo Modul E langkah demi langkah",
        "Ringkasan Modul D, arena, K3, CIS",
    ])

    _add_bullets(prs, "File Penilaian Juri", [
        "Pakai SATU file: templates/LKS2026-Penilaian-Juri.xlsx",
        "File latihan: LKS2026-Penilaian-Juri-CONTOH.xlsx (sudah terisi contoh)",
        "Duplikat per tim (Save As) — jangan pakai file CSV",
        "Kompatibel Excel 2010+ (Microsoft Excel desktop)",
        "Panduan Excel: docs/CONTOH-EXCEL-MODUL-E.md | SOP: docs/SOP-JURI.md",
    ], sub="Marking Scheme CIS = sumber kebenaran poin detail (dari Chief Expert).")

    _add_table(prs, "Cara Isi Excel — Sheet Undian Kubus", [
        "Langkah", "Di mana", "Isi apa",
    ], [
        ["1", "Baris 24 (kuning)", "Warna marker: Merah / Hijau / Biru per slot"],
        ["2", "Baris 28–36 kolom A", "ID kubus: K-M1, K-H2, dst."],
        ["3", "Kolom E", "Rak/Stand: Rak 1, Stand A, ... (dropdown)"],
        ["4", "Kolom F", "Slot: kiri / tengah / kanan (dropdown)"],
        ["5", "Sheet Modul E", "Otomatis — juri isi Hasil Run (1/0) saja"],
    ])

    _add_bullets(prs, "Identitas Lomba", [
        "Bidang: Robot Bergerak Otonom (Autonomous Mobile Robotics)",
        "Tim: 2 peserta | Offline | ±15 jam kerja / 3 hari",
        "Tema 2026: Service Robot at Library",
        "Robot mengambil kubus (buku) dan menaruh ke rak sesuai marker warna",
        "Standar: WorldSkills WSOS Skill 23",
    ])

    _add_table(prs, "Bobot Modul (total 100 poin)", [
        "Modul", "Kriteria", "Hari", "Bobot", "Metode",
    ], [
        ["A", "Organisasi & manajemen kerja", "H1–H3", "8", "Judgement"],
        ["B", "Jurnal / laporan teknis", "H1", "10", "J + M"],
        ["C", "Perakitan robot (L-Channel)", "H1", "10", "Judgement"],
        ["D", "Gerakan dasar & manajemen objek", "H2", "12", "Measurement"],
        ["E", "Performansi otonom", "H2–H3", "60", "Measurement"],
    ])

    _add_bullets(prs, "Sistem Penilaian", [
        "Judgement (J): skala 0–3 untuk proses, kerapian, kolaborasi",
        "Measurement (M): 1 = berhasil, 0 = gagal",
        "3 juri untuk J → rata-rata; selisih ≥ 2 poin → ulangi",
        "Input ke CIS → skala 700 (Medallion jika > 700)",
        "60% skor dari Modul E — fokus utama penilaian juri",
    ])

    _add_table(prs, "Jadwal Singkat", [
        "Hari", "Fokus penilaian juri",
    ], [
        ["H-1", "Technical meeting + familiarisasi arena (2 jam)"],
        ["H1", "Modul A, B, C — perakitan L-Channel (maks. 2 jam)"],
        ["H2", "Modul D (15 menit/tim) + trial Modul E"],
        ["H3", "Final Modul E (+ variasi soal 30%)"],
    ])

    _add_section(prs, "Bagian 2", "Alur Run Otonom — Modul E (60%)")
    _add_flow(prs)

    _add_demo_step(prs, 1, "Briefing layout arena", [
        "Jelaskan posisi START, 3 rak, gate, dan obstacle",
        "Tunjukkan marker warna di slot rak (bisa berubah tiap soal)",
        "Pastikan lampu arena menyala",
        "Isi sheet Undian Kubus di Excel",
    ], "Jangan mulai run sebelum peserta paham layout.")

    _add_demo_step(prs, 2, "Programming & persiapan", [
        "Peserta memprogram navigasi, visi, pick-place",
        "Peserta BOLEH menyentuh laptop",
        "Cek tombol START di robot dan STOP merah",
        "Juri observasi Modul A (kolaborasi & 5S)",
    ])

    _add_demo_step(prs, 3, "Sinyal SIAP", [
        "Peserta menyatakan SIAP saat program siap",
        "Setelah SIAP: TIDAK BOLEH sentuh laptop",
        "Juri pastikan tangan peserta menjauh dari laptop",
        "Siapkan pengacakan kubus",
    ], "Sentuh laptop setelah SIAP = pelanggaran.")

    _add_demo_step(prs, 4, "Pengacakan & isi Excel Undian Kubus", [
        "Acak 9 kubus di rak dan stand (beberapa sengaja salah)",
        "Baris 24: isi warna marker per slot (Merah/Hijau/Biru)",
        "Baris 28–36: kolom A = ID | E = Rak/Stand | F = Slot",
        "Warna, bentuk, lokasi, slot target terisi otomatis",
    ], "Isi Undian SEBELUM run — Modul E membaca data dari sini.")

    _add_demo_step(prs, 5, "Sinyal START & run otonom", [
        "Juri beri sinyal START",
        "Peserta tekan tombol START di robot",
        "Robot berjalan sepenuhnya otonom",
        "Catat di sheet Log Run Otonom",
    ], "Jangan sentuh robot setelah START.")

    _add_demo_step(prs, 6, "Penilaian per kubus (Modul E)", [
        "Sukses = kubus di slot rak, marker warna cocok, stabil",
        "Kolom OK Awal = sudah benar sebelum run (referensi)",
        "Kolom Hasil Run (1/0) = keputusan juri SETELAH run",
        "Bobot: 60/9 poin per kubus (total 60) — konfirmasi ke CE",
    ])

    _add_table(prs, "Contoh baris Modul E (setelah Undian terisi)", [
        "ID", "Posisi Awal", "Slot Target", "OK Awal", "Hasil Run",
    ], [
        ["K-M2", "Stand A — depan", "R1-S1, R2-S2, R3-S3", "0", "1"],
        ["K-H2", "Stand B — depan", "R1-S2, R2-S3, R3-S1", "0", "0"],
        ["K-M1", "Rak 1 — kiri", "R1-S1, R2-S2, R3-S3", "1", "1"],
    ])

    _add_demo_step(prs, 7, "Prosedur RETRY", [
        "Peserta boleh minta RETRY dan ubah program",
        "Kubus di robot dikeluarkan ke luar arena",
        "Posisi kubus tidak dikembalikan ke awal",
        "SIAP lagi → acak ulang → START baru",
    ])

    _add_demo_step(prs, 8, "Rekap & CIS", [
        "Periksa sheet Rekapitulasi di Excel",
        "Chief Expert review pelanggaran & outlier",
        "Input skor resmi ke CIS",
        "Penentuan juara & Medallion of Excellence",
    ])

    _add_bullets(prs, "Latihan bersama juri", [
        "Buka templates/LKS2026-Penilaian-Juri-CONTOH.xlsx (contoh terisi)",
        "Baca docs/CONTOH-EXCEL-MODUL-E.md dan SKENARIO-SOAL-CONTOH.md",
        "Praktikkan: isi Undian → cek Modul E → isi Hasil Run",
        "Simulasikan RETRY dan pelanggaran SIAP/START",
    ])

    _add_section(prs, "Bagian 3", "Ringkasan Modul D, Arena & K3")

    _add_bullets(prs, "Modul D (12 poin, 15 menit/tim)", [
        "Gerak 4 arah, deteksi dinding, line follower U",
        "Deteksi warna: merah, hijau, biru (tampil di monitor)",
        "Deteksi bentuk: O-merah, O-hijau",
        "Pick dari rak & place ke stand — otonom",
    ], sub="Bobot item mentah 12,9 — diskalakan ke 12 di Excel.")

    _add_bullets(prs, "Arena & spesifikasi robot", [
        "Arena ±4040×2040 mm | 3 rak | 9 kubus | 10 obstacle | 1 gate",
        "Controller: Raspberry Pi atau myRIO",
        "Python / C++ / LabVIEW | baterai ≤12V",
        "Tombol STOP merah wajib | 30% rule di H3",
    ])

    _add_bullets(prs, "K3 & klarifikasi ke Chief Expert", [
        "Larangan air terbuka dekat peralatan listrik",
        "Reset obstacle yang jatuh sebelum tim berikutnya",
        "Klarifikasi: ukuran arena, max RETRY, bobot CIS",
        "Dokumen resmi: smk.pusatprestasinasional.kemdikdasmen.go.id",
    ], sub="Terima kasih — selamat bertugas sebagai juri LKS 2026!")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Created: {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
